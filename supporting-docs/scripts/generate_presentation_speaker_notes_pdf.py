#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_BREAK
from docx.shared import Cm, Pt
from pptx import Presentation


NOTES = [
    {
        "key_point": "Open with a clean one-sentence definition of the project so the audience knows immediately that this is a benchmark, not a composition system.",
        "talk_track": (
            "This project proposes a symbolic music benchmark for music theory and computational musicology. "
            "The core idea is to evaluate whether language models can solve structured music-theory tasks, such as identifying intervals or transforming melodies, using automatically checkable outputs. "
            "The benchmark combines classification tasks, transformation tasks, formal evaluation, and an optional analytical explanation layer for paper-oriented analysis."
        ),
        "likely_question": "Is this a music-generation tool or a benchmark?",
        "safe_answer": "It is primarily a benchmark. The goal is to measure music-theoretical reasoning in symbolic form, not to generate finished musical pieces.",
        "transition": "After defining the project, I move to why symbolic benchmarking is a useful starting point.",
    },
    {
        "key_point": "Use this slide only to orient the room and set expectations for the flow.",
        "talk_track": (
            "I will first explain the motivation and research question, then the benchmark tasks and scope, then the pipeline and schema, then how evaluation works, and finally why this can also be framed as a musicological question. "
            "So the talk moves from why this benchmark matters, to how it is built, to what kind of research it enables."
        ),
        "likely_question": "Are all sections equally important for the proposal?",
        "safe_answer": "No. The most important pieces are the motivation, benchmark design, evaluation logic, and the musicological framing.",
        "transition": "With the structure in mind, the next slide explains why we chose a symbolic benchmark in the first place.",
    },
    {
        "key_point": "Frame symbolic evaluation as a design choice for reproducibility and clarity, not as a claim that audio is unimportant.",
        "talk_track": (
            "We chose a symbolic benchmark because symbolic representations let us define music-theoretical relationships that can be checked by rules. "
            "That makes the benchmark reproducible and easier to evaluate automatically than an audio-first setup. "
            "It also means we can compare models on musically meaningful structure rather than only on surface fluency. "
            "So the benchmark is designed to be simple, testable, and usable even by non-expert users."
        ),
        "likely_question": "Why not work directly with audio?",
        "safe_answer": "Audio is important, but for a first benchmark release we wanted a controlled setting with stable ground truth and transparent evaluation. Symbolic music is the cleanest place to start that process.",
        "transition": "Once that motivation is clear, the next step is the actual research question behind the benchmark.",
    },
    {
        "key_point": "State the research question clearly and emphasize that the second half is about the gap between correctness and understanding.",
        "talk_track": (
            "Our main question is whether large language models can capture and manipulate core music-theoretical structures in symbolic form. "
            "But there is a second question that matters just as much: even if a model gives the correct symbolic answer, does that reflect meaningful analytical reasoning, or just pattern matching? "
            "The benchmark is designed to help us study both structural correctness and the limits of explanation."
        ),
        "likely_question": "Do you claim this measures true musical understanding?",
        "safe_answer": "Not directly. What we claim is that it probes structured music-theoretical competence and helps reveal where formal correctness and analytical reasoning diverge.",
        "transition": "To position that question, the next slide briefly places the work against related symbolic music research.",
    },
    {
        "key_point": "Present the novelty as evaluation focus, not as replacing all prior symbolic music work.",
        "talk_track": (
            "There is already strong work on symbolic music representations and music generation, including token-based approaches such as REMI-like encodings. "
            "However, much of that literature focuses on representation quality or generation quality, not on a compact benchmark for music-theoretical reasoning. "
            "Our contribution is to operationalize theory concepts as tasks with stable inputs, stable outputs, and automatic evaluation."
        ),
        "likely_question": "What is new compared with existing symbolic music models?",
        "safe_answer": "The novelty is not a new representation model by itself. It is the benchmark design: task formulation, evaluation protocol, and the bridge to musicological interpretation.",
        "transition": "After that positioning, I can describe the actual task families in the benchmark.",
    },
    {
        "key_point": "Make the benchmark feel concrete by splitting it into two intuitive task families.",
        "talk_track": (
            "The benchmark has eight tasks grouped into two families. "
            "The first family is label classification, where the model predicts a label such as an interval, a chord type, a harmonic function, or a voice-leading violation. "
            "The second family is sequence transformation, where the model returns a transformed note sequence, for example under transposition, inversion, retrograde, or rhythmic scaling. "
            "This lets us evaluate both recognition and manipulation of musical structure."
        ),
        "likely_question": "Why these eight tasks specifically?",
        "safe_answer": "Because together they cover a useful spread of theory skills: identifying relations, classifying harmonic meaning, detecting rule violations, and applying explicit transformations.",
        "transition": "The next slide clarifies what is intentionally included in version one and what is intentionally left narrow.",
    },
    {
        "key_point": "Defend the narrow scope as an intentional V1 boundary that makes the benchmark robust.",
        "talk_track": (
            "Version one is intentionally narrow. Harmonic function is limited to major mode and three broad labels. "
            "Voice-leading is simplified to two voices with a small label set. "
            "Harmony uses a reduced chord vocabulary, and the benchmark stays fully symbolic rather than audio-based. "
            "The point is not to oversimplify music forever, but to define a stable first slice that can be evaluated cleanly."
        ),
        "likely_question": "Isn't this too limited to support a serious paper?",
        "safe_answer": "It is limited, but intentionally so. The scientific value comes from having a controlled, reproducible benchmark first, then extending it in future iterations.",
        "transition": "That is why the next slide makes the planned extensions explicit instead of hiding them.",
    },
    {
        "key_point": "Show that V1 is usable now, but not framed as final theoretical coverage.",
        "talk_track": (
            "This slide makes an important methodological point: the current benchmark is already scientifically usable, but it should not be presented as the final theoretical boundary of the paper. "
            "We already know where natural extensions lie, including richer harmonic-function coverage, broader voice-leading analysis, expanded harmonic vocabulary, and more developed explanation analysis. "
            "So V1 is the current benchmark slice, not the final horizon of the project."
        ),
        "likely_question": "Why present the benchmark before all extensions are implemented?",
        "safe_answer": "Because the current version is already coherent, runnable, and evaluable. That makes it a valid first benchmark release and a solid base for future work.",
        "transition": "With scope established, I can now show the technical pipeline that implements the benchmark.",
    },
    {
        "key_point": "Present the architecture as a pipeline from rule-based theory logic to model evaluation outputs.",
        "talk_track": (
            "The benchmark pipeline starts from theory rules and task specifications, which generate controlled symbolic benchmark instances. "
            "Those instances are exported into tokenizer-specific views and model-ready cases. "
            "Then models are run, predictions are collected, and a unified evaluator produces overall summaries and analysis artifacts. "
            "So the architecture separates theory generation, representation, execution, and analysis."
        ),
        "likely_question": "Are the benchmark cases synthetic?",
        "safe_answer": "Yes, they are rule-based generated cases. That is a feature here, because it gives us controlled ground truth and consistent evaluation.",
        "transition": "The next slide turns that architecture into the simpler user-facing workflow.",
    },
    {
        "key_point": "Emphasize usability for someone who is not deeply familiar with the repository.",
        "talk_track": (
            "From a user perspective, the workflow is straightforward: generate the benchmark cases, export representation views, validate them, build model input files, run a model, and evaluate the outputs. "
            "That standard path is meant to be easy to document and easy to execute end to end. "
            "In other words, the benchmark is not only theoretically motivated, but also packaged to be practically usable."
        ),
        "likely_question": "How difficult is it to run this benchmark?",
        "safe_answer": "The standard path is intentionally simple. A user can follow the scripted workflow without needing to understand every internal module in the repository.",
        "transition": "After the runtime flow, the next slide defines the output schema that keeps the benchmark stable across tasks.",
    },
    {
        "key_point": "Explain the schema as the contract that keeps evaluation stable while preserving backward compatibility.",
        "talk_track": (
            "The canonical prediction schema distinguishes clearly between label tasks and note-sequence tasks. "
            "For label tasks, the main field is prediction_label, and optionally a model can add prediction_explanation. "
            "For sequence tasks, the main field is prediction_notes, with optional structured decoding. "
            "At the same time, we keep the legacy prediction field for backward compatibility, so older tooling does not break."
        ),
        "likely_question": "Why keep the older prediction field if you already defined a better schema?",
        "safe_answer": "Because compatibility matters. The new schema improves clarity, but the legacy field lets us migrate existing tools and outputs smoothly.",
        "transition": "Once the schema is defined, we can explain how scoring works across all eight tasks.",
    },
    {
        "key_point": "Stress that evaluation is unified overall but still sensitive to the different structure of each task.",
        "talk_track": (
            "The evaluator covers all eight tasks in one standard pipeline, but it still uses task-appropriate metrics. "
            "Classification tasks use exact match, accuracy, and macro F1 where appropriate, while sequence tasks use measures such as exact match, pitch accuracy, interval preservation, rhythm preservation, or bar validity. "
            "This gives us a single benchmark framework without pretending that every music task should be evaluated in exactly the same way."
        ),
        "likely_question": "Why not use one universal metric for everything?",
        "safe_answer": "Because the structure of the tasks is genuinely different. A label classification task and a transformed melody require different notions of correctness.",
        "transition": "The next three slides show example result views and how the benchmark differentiates between models and tokenizers.",
    },
    {
        "key_point": "Introduce the first result view as a tokenizer comparison within model runs, not as a final leaderboard claim.",
        "talk_track": (
            "This first results view compares tokenizers task by task within specific model runs. "
            "The main point is not to claim a definitive ranking, but to show that the benchmark is sensitive enough to capture meaningful differences across representations and tasks. "
            "Some tasks appear relatively robust, while others vary more across tokenizers, which suggests that representation design matters for symbolic music reasoning."
        ),
        "likely_question": "Are these final benchmark results?",
        "safe_answer": "They should be treated as current benchmark runs that validate the framework and illustrate model-tokenizer differences, not as the final last word on performance.",
        "transition": "The next view flips the perspective and compares models under particular tokenizer settings.",
    },
    {
        "key_point": "Use this slide to talk about consistency across models rather than reading every bar.",
        "talk_track": (
            "Here the comparison is organized by tokenizer view, so we can see how different models behave under the same representation. "
            "The broad pattern is that stronger models tend to be more consistent across tasks, but difficult tasks still remain difficult across the board. "
            "So this is less about one single winning bar and more about showing that the benchmark discriminates meaningfully across model families."
        ),
        "likely_question": "What is the main takeaway from this chart?",
        "safe_answer": "The benchmark is informative: it captures both cross-model differences and task difficulty patterns, instead of collapsing everything into a trivial score.",
        "transition": "The last result view aggregates performance more globally across all tokenizer configurations.",
    },
    {
        "key_point": "Summarize the aggregate picture without overexplaining the chart.",
        "talk_track": (
            "This aggregated view gives the highest-level picture across models and tokenizer configurations. "
            "The value of this slide is that it makes it easier to see which tasks are systematically easier, which remain challenging, and how performance changes when we collapse across representation choices. "
            "Again, the purpose here is not only to report numbers, but to show that the benchmark produces structured and interpretable differences."
        ),
        "likely_question": "What should the audience remember from the evaluation slides overall?",
        "safe_answer": "The key point is that the benchmark is not flat. It can distinguish between tasks, models, and tokenizers in a way that supports analysis.",
        "transition": "After the formal evaluation layer, I move to the extra musicological layer of the benchmark.",
    },
    {
        "key_point": "Present standard mode and explanatory mode as two layers with different goals.",
        "talk_track": (
            "This distinction is important for both usability and scientific framing. "
            "The standard benchmark mode is the mandatory path: it gives stable outputs, easy parsing, and automatic scoring, which is ideal for demos, testing, and reproducibility. "
            "The explanatory mode is optional and paper-oriented: it keeps the same main outputs but may add prediction_explanation so that we can compare structural correctness with analytical language. "
            "This way, explanation enriches the project without destabilizing the core benchmark."
        ),
        "likely_question": "Why not make explanations mandatory?",
        "safe_answer": "Because explanation quality is harder to score reliably and would make the core benchmark less stable. We want the main benchmark to stay simple and reproducible.",
        "transition": "That leads directly to the main musicological question behind the project.",
    },
    {
        "key_point": "Use this slide to articulate the musicological value in one clear sentence.",
        "talk_track": (
            "The benchmark is testing more than pattern recall. It probes harmonic function, voice-leading constraints, and structural transformations, which are all meaningful music-theoretical categories. "
            "So the musicological question is whether correct symbolic answers actually reflect musical understanding, or whether models are only matching token patterns. "
            "Our expected outcome is to identify cases where correctness and reasoning align, and cases where they clearly diverge."
        ),
        "likely_question": "How is this musicology rather than only AI evaluation?",
        "safe_answer": "Because the benchmark operationalizes analytical categories from music theory and uses them to study how models handle musical structure and explanation.",
        "transition": "The next few slides show why this benchmark can serve different users and how its internal logic is organized.",
    },
    {
        "key_point": "Explain the use-case diagram as a communication tool, not as the central scientific result.",
        "talk_track": (
            "This use-case view helps explain the benchmark to different audiences. "
            "A non-technical user mainly needs a simple standard path to run and score the benchmark. "
            "A research user may focus on annotated outputs and qualitative examples for the paper. "
            "A model developer may use the benchmark to compare models or tokenizer choices. "
            "So the same benchmark supports practical evaluation, analysis, and experimentation."
        ),
        "likely_question": "Who is the primary audience of the benchmark?",
        "safe_answer": "The benchmark is designed to be accessible to non-expert users, but it also supports research analysis and model comparison.",
        "transition": "After that user-oriented view, the next slide shows the runtime order of the benchmark components.",
    },
    {
        "key_point": "Use the sequence diagram to make the execution order intuitive.",
        "talk_track": (
            "This sequence view simply shows the order of operations in the benchmark pipeline. "
            "A user triggers generation, then export and case building, then model execution, then evaluation, and finally review of outputs. "
            "It is useful because it separates the temporal flow of the system from the logical grouping of modules."
        ),
        "likely_question": "Can these steps be modified or swapped?",
        "safe_answer": "Yes. The pipeline is modular enough that models, tokenizers, and some analysis steps can be changed without redesigning the whole benchmark.",
        "transition": "The next slide complements runtime order with a higher-level component organization.",
    },
    {
        "key_point": "Highlight separation of concerns across theory logic, representation, execution, and analysis.",
        "talk_track": (
            "This component view groups the repository into four main areas. "
            "Theory logic defines the benchmark rules and task semantics. "
            "Representation handles tokenizers and prediction formats. "
            "Execution builds cases and runs models. "
            "Analysis evaluates outputs and produces review artifacts. "
            "That separation makes the benchmark easier to maintain and easier to extend."
        ),
        "likely_question": "Where would you add future extensions first?",
        "safe_answer": "Most likely in the theory and analysis layers: richer task rules, broader coverage, and more detailed explanation analysis.",
        "transition": "With the system organization covered, I can now summarize the proposal in one slide.",
    },
    {
        "key_point": "Reduce this slide to the single strongest takeaway of the whole presentation.",
        "talk_track": (
            "The summary is that this project investigates whether language models can correctly capture and manipulate core music-theoretical structure in symbolic form, and whether that structural correctness aligns with analytical explanation quality. "
            "So the benchmark is meant to measure both performance and the limits of performance as evidence of musical reasoning."
        ),
        "likely_question": "What is the one-sentence contribution?",
        "safe_answer": "A reproducible multi-task symbolic music-theory benchmark with automatic evaluation and an optional analytical explanation layer.",
        "transition": "The final slide frames that contribution in paper terms and points to the immediate next steps.",
    },
    {
        "key_point": "Close by tying benchmark design, paper contribution, and next steps together.",
        "talk_track": (
            "In paper terms, the contribution is a symbolic benchmark for music theory and computational musicology, a multi-task setup for musical reasoning, and an evaluation framework that separates stable benchmarking from optional explanation analysis. "
            "The current benchmark is a V1 scope with planned extensions, not final theoretical coverage. "
            "The immediate next steps are to strengthen API-based experimentation, run broader benchmark experiments, and collect explanation examples that help us analyze the gap between correctness and reasoning."
        ),
        "likely_question": "What is the next concrete milestone after this proposal?",
        "safe_answer": "The next milestone is a larger experimental run with multiple models and tokenizers, followed by qualitative analysis of explanatory cases for the paper.",
        "transition": "End by thanking the audience and invite questions on the benchmark design, the evaluation, or the musicological framing.",
    },
]


GENERAL_QA = [
    (
        "Why symbolic music instead of audio?",
        "Because symbolic music gives a controllable first benchmark with explicit structure and reliable ground truth. Audio can be a later extension once the evaluation logic is stable.",
    ),
    (
        "Why is major mode the only harmonic-function setting in V1?",
        "Because the first goal is reproducibility. Major mode provides a narrow but stable starting point for task definition and automatic verification.",
    ),
    (
        "Why only two voices for voice-leading?",
        "Two voices let us define clean rule checks first. It is a deliberate simplification, not the final intended scope of the project.",
    ),
    (
        "How is this different from ordinary prompt testing?",
        "The benchmark provides controlled inputs, canonical outputs, rule-based gold references, and a unified evaluator. That makes it a benchmark, not just a collection of prompts.",
    ),
    (
        "How does this respond to a musicological question?",
        "It operationalizes analytical categories such as harmonic function and voice-leading, then asks whether formal correctness aligns with analytical reasoning. That links benchmark performance to a real interpretive question.",
    ),
    (
        "Why not score explanations automatically as part of the main benchmark?",
        "Because explanation quality is more subjective and harder to annotate reliably. Keeping explanations optional protects the robustness of the main benchmark.",
    ),
    (
        "Are the accuracy plots final results?",
        "They should be presented as current benchmark runs that demonstrate the benchmark is sensitive to model and tokenizer differences. The core contribution is the framework and evaluation protocol.",
    ),
    (
        "What would count as success for this project?",
        "A successful outcome is a reproducible benchmark that clearly differentiates tasks, models, and tokenizers, and that helps analyze the gap between structural correctness and analytical explanation.",
    ),
]


def first_non_empty_text(slide) -> str:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                return " ".join(text.split())
    return "Untitled slide"


def slide_image_path(rendered_dir: Path, slide_index: int) -> Path | None:
    path = rendered_dir / f"slide-{slide_index:02d}.png"
    return path if path.exists() else None


def add_label_paragraph(doc: Document, label: str, body: str) -> None:
    para = doc.add_paragraph()
    run = para.add_run(f"{label}: ")
    run.bold = True
    para.add_run(body)


def write_markdown(md_path: Path, slide_titles: list[str]) -> None:
    lines = [
        "# Speaker Notes",
        "",
        "These notes are designed as a tablet-friendly companion for the presentation.",
        "Each slide includes a short talk track, one likely question, and a safe answer.",
        "",
        "## Suggested pacing",
        "",
        "- Slides 1-5: motivation and framing",
        "- Slides 6-12: benchmark design and evaluation logic",
        "- Slides 13-15: results overview, speak comparatively and do not overclaim",
        "- Slides 16-20: musicology and system views",
        "- Slides 21-22: summary and next steps",
        "",
    ]

    for idx, (title, notes) in enumerate(zip(slide_titles, NOTES), start=1):
        lines.extend(
            [
                f"## Slide {idx}: {title}",
                "",
                f"**Key point:** {notes['key_point']}",
                "",
                f"**Talk track:** {notes['talk_track']}",
                "",
                f"**Likely question:** {notes['likely_question']}",
                "",
                f"**Safe answer:** {notes['safe_answer']}",
                "",
                f"**Transition:** {notes['transition']}",
                "",
            ]
        )

    lines.extend(["## General Q&A", ""])
    for question, answer in GENERAL_QA:
        lines.extend([f"**Q:** {question}", "", f"**A:** {answer}", ""])

    md_path.write_text("\n".join(lines), encoding="utf-8")


def build_docx(pptx_path: Path, rendered_dir: Path, docx_path: Path, md_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    slide_titles = [first_non_empty_text(slide) for slide in prs.slides]

    if len(slide_titles) != len(NOTES):
        raise ValueError(
            f"Expected {len(NOTES)} slides for notes generation, found {len(slide_titles)} slides."
        )

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(1.4)
    section.bottom_margin = Cm(1.4)
    section.left_margin = Cm(1.6)
    section.right_margin = Cm(1.6)

    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Aptos"
    normal_style.font.size = Pt(10.5)

    title_style = doc.styles["Title"]
    title_style.font.name = "Aptos Display"
    title_style.font.size = Pt(22)

    heading1 = doc.styles["Heading 1"]
    heading1.font.name = "Aptos Display"
    heading1.font.size = Pt(16)

    heading2 = doc.styles["Heading 2"]
    heading2.font.name = "Aptos Display"
    heading2.font.size = Pt(12)

    doc.add_heading("Speaker Notes", 0)
    intro = doc.add_paragraph()
    intro.add_run(
        "Tablet-friendly notes for the presentation "
    ).bold = True
    intro.add_run(
        f"\"{pptx_path.stem}\". Each slide includes a short talk track, one likely question, and a safe response."
    )

    pacing = doc.add_paragraph(style="List Bullet")
    pacing.add_run("Slides 1-5: motivation and framing")
    for line in [
        "Slides 6-12: benchmark design, pipeline, schema, and evaluation",
        "Slides 13-15: results overview, speak comparatively and avoid overclaiming",
        "Slides 16-20: musicology and system views",
        "Slides 21-22: summary and next steps",
    ]:
        para = doc.add_paragraph(style="List Bullet")
        para.add_run(line)

    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)

    for idx, (title, notes) in enumerate(zip(slide_titles, NOTES), start=1):
        doc.add_heading(f"Slide {idx}. {title}", level=1)
        image_path = slide_image_path(rendered_dir, idx)
        if image_path:
            doc.add_picture(str(image_path), width=Cm(16.7))

        add_label_paragraph(doc, "Key point", notes["key_point"])
        add_label_paragraph(doc, "Suggested talk track", notes["talk_track"])
        add_label_paragraph(doc, "Likely question", notes["likely_question"])
        add_label_paragraph(doc, "Safe answer", notes["safe_answer"])
        add_label_paragraph(doc, "Transition", notes["transition"])

        if idx != len(NOTES):
            doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)

    doc.add_section(WD_SECTION.NEW_PAGE)
    doc.add_heading("General Q&A", level=1)
    doc.add_paragraph(
        "These are safe short answers for likely forum or professor questions that may come up across multiple slides."
    )
    for question, answer in GENERAL_QA:
        add_label_paragraph(doc, "Question", question)
        add_label_paragraph(doc, "Answer", answer)
        doc.add_paragraph()

    doc.save(str(docx_path))
    write_markdown(md_path, slide_titles)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a presentation speaker-notes DOCX and Markdown.")
    parser.add_argument("--pptx", required=True, help="Path to the source PPTX file.")
    parser.add_argument("--rendered-dir", required=True, help="Directory containing slide PNGs.")
    parser.add_argument("--docx-out", required=True, help="Path for the output DOCX.")
    parser.add_argument("--md-out", required=True, help="Path for the output Markdown.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx).expanduser().resolve()
    rendered_dir = Path(args.rendered_dir).expanduser().resolve()
    docx_path = Path(args.docx_out).expanduser().resolve()
    md_path = Path(args.md_out).expanduser().resolve()

    docx_path.parent.mkdir(parents=True, exist_ok=True)
    md_path.parent.mkdir(parents=True, exist_ok=True)

    build_docx(pptx_path, rendered_dir, docx_path, md_path)
    print(f"Wrote DOCX to {docx_path}")
    print(f"Wrote Markdown to {md_path}")


if __name__ == "__main__":
    main()
